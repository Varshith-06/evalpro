"""Build the SIH 2026 idea-submission deck from the official template.

The deck is generated **from the provided template file itself** rather than
from a blank presentation, so the slide master, layouts, theme, colour scheme,
the SIH banner artwork, the footer bar, the team-name oval, and the slide
dimensions are byte-identical to what the portal supplied. Only text content and
the geometry of the content text boxes change.

Format decisions, and the reasoning:

* **Six slides.** The template ships seven and its own final slide says to keep
  the deck to six including the title page, and that the instructions slide may
  be deleted before upload. So it is deleted.
* **The idea-details pointers are preserved.** The template forbids changing
  them, so every pointer ("Detailed explanation of the proposed solution",
  "Potential challenges and risks", and so on) stays, verbatim, as the bold
  underlined sub-heading it already is. The team's content goes underneath it.
* **Typefaces are the template's.** Garamond for the SIH banner, Times New Roman
  for slide titles, Arial for all body text, with the template's Wingdings and
  Arial bullet characters and its ``tx2`` heading colour.
* **Body point sizes are reduced from the template's 28pt placeholders.** Real
  content does not fit at 28pt in a box 1.5 inches tall. The typeface, colour,
  bullet glyphs, alignment, and every other formatting property are unchanged.

Run with:  python scripts/build_presentation.py
"""
from __future__ import annotations

import copy
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parent.parent
TEMPLATE = REPO / "SIH2026-IDEA-Presentation-Format.pptx"
OUTPUT = REPO / "presentation" / "SIH2026-EvalPro-Idea-Presentation.pptx"

# Template palette, read from the file rather than guessed.
HEADING_BLUE = RGBColor(0x1F, 0x49, 0x7D)   # theme tx2
BODY_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT_BLUE = RGBColor(0x00, 0x70, 0xC0)    # the footer bar colour
LAYER_FILL = RGBColor(0x1F, 0x49, 0x7D)
LAYER_FILL_ALT = RGBColor(0x4F, 0x81, 0xBD)  # theme accent1
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

ARIAL = "Arial"
TIMES = "Times New Roman"


# ==========================================================================
# Low-level helpers
# ==========================================================================
def set_bullet(paragraph, char: str | None, font: str = ARIAL, indent_in: float = 0.0) -> None:
    """Set (or clear) a paragraph's bullet glyph.

    python-pptx has no API for this, so it is done on the XML directly using
    exactly the ``buFont``/``buChar`` pairs the template already uses.
    """
    p_pr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont"):
        for node in p_pr.findall(qn(tag)):
            p_pr.remove(node)

    marl = int(Inches(indent_in + (0.24 if char else 0.0)))
    p_pr.set("marL", str(marl))
    p_pr.set("indent", str(-int(Inches(0.24)) if char else "0"))

    if char is None:
        p_pr.append(p_pr.makeelement(qn("a:buNone"), {}))
        return

    bu_font = p_pr.makeelement(qn("a:buFont"), {"typeface": font, "pitchFamily": "34", "charset": "0"})
    bu_char = p_pr.makeelement(qn("a:buChar"), {"char": char})
    p_pr.append(bu_font)
    p_pr.append(bu_char)


def set_line_spacing(paragraph, spacing: float, before: float = 0.0, after: float = 0.0) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    for tag in ("a:lnSpc", "a:spcBef", "a:spcAft"):
        for node in p_pr.findall(qn(tag)):
            p_pr.remove(node)
    ln = p_pr.makeelement(qn("a:lnSpc"), {})
    pct = ln.makeelement(qn("a:spcPct"), {"val": str(int(spacing * 100000))})
    ln.append(pct)
    p_pr.insert(0, ln)
    for tag, value in (("a:spcBef", before), ("a:spcAft", after)):
        if value <= 0:
            continue
        node = p_pr.makeelement(qn(tag), {})
        pts = node.makeelement(qn("a:spcPts"), {"val": str(int(value * 100))})
        node.append(pts)
        p_pr.append(node)


def style_run(run, size: float, bold: bool = False, colour: RGBColor = BODY_BLACK,
              font: str = ARIAL, underline: bool = False, italic: bool = False) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.underline = underline
    run.font.color.rgb = colour
    # Match the template, which sets the complex-script typeface alongside latin.
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:cs", "a:ea"):
        for node in r_pr.findall(qn(tag)):
            r_pr.remove(node)
    r_pr.append(r_pr.makeelement(qn("a:cs"), {"typeface": font, "pitchFamily": "34", "charset": "0"}))


def clear_text(shape) -> None:
    frame = shape.text_frame
    frame.clear()
    first = frame.paragraphs[0]
    for run in list(first.runs):
        first._p.remove(run._r)


def add_line(frame, text: str, *, size: float, bold: bool = False,
             colour: RGBColor = BODY_BLACK, bullet: str | None = None,
             indent: float = 0.0, underline: bool = False, spacing: float = 1.0,
             before: float = 0.0, after: float = 0.0, first: bool = False,
             align=PP_ALIGN.LEFT, font: str = ARIAL, bullet_font: str = ARIAL):
    paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    style_run(run, size, bold=bold, colour=colour, underline=underline, font=font)
    set_bullet(paragraph, bullet, font=bullet_font, indent_in=indent)
    set_line_spacing(paragraph, spacing, before, after)
    return paragraph


def drop_empty_paragraphs(frame) -> None:
    """Remove paragraphs with no runs.

    Some of the template's boxes carry a list style that suppresses the bullet
    on the very first paragraph. Building every paragraph the same way and then
    dropping the leftover empty one is more reliable than special-casing which
    box behaves how.
    """
    body = frame._txBody
    for paragraph in list(frame.paragraphs):
        if not paragraph.runs:
            body.remove(paragraph._p)
    if not frame.paragraphs:
        body.append(body.makeelement(qn("a:p"), {}))


def place(shape, left: float, top: float, width: float, height: float) -> None:
    shape.left, shape.top, shape.width, shape.height = (
        Inches(left), Inches(top), Inches(width), Inches(height),
    )


def find(slide, name_fragment: str):
    for shape in slide.shapes:
        if name_fragment.lower() in (shape.name or "").lower():
            return shape
    raise KeyError(f"{name_fragment!r} not found on slide")


def set_autofit_off(shape) -> None:
    """The template's boxes carry ``spAutoFit``; with denser content we size the
    box ourselves, so the tag is removed rather than fought with."""
    body = shape.text_frame._txBody.find(qn("a:bodyPr"))
    if body is None:
        return
    for tag in ("a:spAutoFit", "a:normAutofit"):
        for node in body.findall(qn(tag)):
            body.remove(node)
    body.set("wrap", "square")
    for attribute, value in (("lIns", "45720"), ("rIns", "45720"), ("tIns", "22860"), ("bIns", "22860")):
        body.set(attribute, value)


def delete_slide(prs, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    entries = list(slide_id_list)
    entry = entries[index]
    prs.part.drop_rel(entry.rId)
    slide_id_list.remove(entry)


# ==========================================================================
# Slide content
# ==========================================================================
TEAM = {
    # Fields the portal issues. Left as obvious placeholders to be filled in
    # before upload rather than invented here.
    "ps_id": "[Problem Statement ID from the SIH portal]",
    "ps_title": "Automated Programming Lab Evaluation Platform",
    "theme": "Smart Education",
    "category": "Software",
    "team_id": "[Team ID from the SIH portal]",
    "team_name": "[Team Name as registered on the portal]",
}

IDEA_TITLE = "EvalPro: Concept-Level Mastery Analytics for Programming Labs"


def build_title_slide(slide) -> None:
    box = find(slide, "TextBox 9")
    set_autofit_off(box)
    place(box, 0.36, 2.10, 6.20, 5.05)
    clear_text(box)
    frame = box.text_frame

    fields = [
        ("Problem Statement ID – ", TEAM["ps_id"]),
        ("Problem Statement Title – ", TEAM["ps_title"]),
        ("Theme – ", TEAM["theme"]),
        ("PS Category – ", TEAM["category"]),
        ("Team ID – ", TEAM["team_id"]),
        ("Team Name (Registered on portal) – ", TEAM["team_name"]),
    ]
    for index, (label, value) in enumerate(fields):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.JUSTIFY
        label_run = paragraph.add_run()
        label_run.text = label
        style_run(label_run, 15, bold=True, colour=BODY_BLACK)
        value_run = paragraph.add_run()
        value_run.text = value
        style_run(value_run, 15, bold=False, colour=HEADING_BLUE)
        set_bullet(paragraph, "•")
        set_line_spacing(paragraph, 1.15, before=0, after=12)

    # The template's own subtitle placeholder reads "TITLE PAGE"; it stays.
    subtitle = find(slide, "Subtitle 3")
    for paragraph in subtitle.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip().upper() == "TITLE PAGE":
                run.text = "TITLE PAGE"


def build_solution_slide(slide) -> None:
    title = find(slide, "Title 1")
    clear_text(title)
    add_line(
        title.text_frame, IDEA_TITLE, size=26, bold=True, colour=HEADING_BLUE,
        font=TIMES, first=True, align=PP_ALIGN.LEFT,
    )
    # The template's team-name oval sits at x 0.36–1.73, so the title starts
    # clear of it rather than running underneath.
    place(title, 1.88, 0.04, 8.55, 1.08)

    box = find(slide, "TextBox 8")
    set_autofit_off(box)
    place(box, 0.36, 1.24, 12.58, 4.70)
    clear_text(box)
    frame = box.text_frame

    add_line(
        frame,
        "Proposed Solution (Describe your Idea/Solution/Prototype)",
        size=17, bold=True, colour=HEADING_BLUE, underline=True,
        bullet="v", bullet_font="Wingdings", first=True, after=7,
    )

    sections = [
        (
            "Detailed explanation of the proposed solution",
            [
                "Not a grader that reports analytics — an analytics platform whose sensor is an automated grader.",
                "The spine is a course concept graph (40–80 nodes, prerequisite DAG, ~2 hours to author once, reused every semester). Every rubric item carries concept_ids, so each submission becomes evidence about named competencies.",
                "A seven-stage cascade grades it: ingest → tree-sitter parse → integrity screen → sandboxed build → sandboxed test → partial credit → confidence gate. Each stage writes evidence; only the gate writes a score.",
                "Bayesian Knowledge Tracing accumulates that evidence into per-student mastery with explicit uncertainty, propagated down the prerequisite DAG.",
            ],
        ),
        (
            "How it addresses the problem",
            [
                "Collect academic information — LTI 1.3 roster and SIS sync, course outcomes, submissions, reports, gradebook writeback.",
                "Analyse student and course data — per submission, per student longitudinally, per rubric item (classical psychometrics), and per cohort.",
                "Actionable insights — prerequisite-walk remediation for students, re-teach signals ranked by downstream impact for faculty, early warning routed to support for administrators.",
                "Simple interface — one landing question per role: what should I work on / what should I teach / where is this programme weak.",
            ],
        ),
        (
            "Innovation and uniqueness of the solution",
            [
                "Partial credit by repair distance: the smallest edit that makes the code compile. A missing colon costs two marks, not a hundred percent of them.",
                "The test oracle never enters the sandbox — student code cannot hardcode an answer it was never given.",
                "Item analysis grades the assessment, not just the student: negative discrimination exposes an ambiguous spec before sixty students hit it.",
                "Abstain rather than guess — a tunable auto-release dial, and an integrity screen that reports cohort-relative outliers as evidence, never as a verdict.",
                "CO–PO attainment computed live from real performance evidence, not reconstructed from a spreadsheet.",
            ],
        ),
    ]
    _render_sections(frame, sections, heading_size=14.5, body_size=11.5)

    metric_strip(slide, [
        ("91", "submissions graded end to end by the real cascade, in the real sandbox"),
        ("24", "concepts tracked per student, mapped to 5 course outcomes"),
        ("0.9 s", "p95 latency for the whole seven-stage cascade, per submission"),
        ("13", "isolation layers specified; the test oracle never enters the guest"),
    ])


def build_technical_slide(slide) -> None:
    box = find(slide, "TextBox 8")
    set_autofit_off(box)
    place(box, 0.36, 3.26, 12.58, 3.52)
    clear_text(box)
    frame = box.text_frame

    add_line(
        frame,
        "Technologies to be used (e.g. programming languages, frameworks, hardware)",
        size=14, bold=True, colour=HEADING_BLUE, underline=True,
        bullet="v", bullet_font="Wingdings", first=True, after=4,
    )
    for text in [
        "Backend — Python 3.12, FastAPI, SQLAlchemy, PostgreSQL (SQLite for the pilot). Build-free ES-module front end with hand-authored SVG, so a deployment needs no toolchain.",
        "Analysis — tree-sitter (40+ grammars, error-tolerant), winnowing fingerprints (MOSS), pq-gram/APTED tree edit distance, Bayesian Knowledge Tracing, HDBSCAN, point-biserial item analysis.",
        "Models — gradient boosting for language ID, confidence estimation and risk; graph embeddings for structural clones; a fine-tuned DeBERTa-v3 encoder for report–code entailment. The LLM drafts rubrics once per assignment and never grades a submission.",
        "Isolation — Firecracker microVM, seccomp-bpf allowlist, cgroups v2, empty network namespace, one-shot instances, supervisor-enforced wall clock.",
        "Integration — LTI 1.3 tool provider with Assignment and Grade Services writeback (Moodle, Canvas, Google Classroom).",
    ]:
        add_line(frame, text, size=11, bullet="•", indent=0.12, spacing=0.95, after=2,
                 align=PP_ALIGN.JUSTIFY)

    add_line(
        frame,
        "Methodology and process for implementation (Flow Charts/Images/ working prototype)",
        size=14, bold=True, colour=HEADING_BLUE, underline=True,
        bullet="v", bullet_font="Wingdings", before=8, after=4,
    )
    for text in [
        "Working prototype: the four layers above run today. A 24-student cohort across four labs — 91 submissions — is graded end to end by the real cascade in a real sandbox in under 90 seconds on a laptop; p95 latency is 0.9 s per submission.",
        "Phase 1 trustworthy execution → Phase 2 partial credit and authoring → Phase 3 accumulation → Phase 4 action → Phase 5 semantic depth. Each phase ships something separately useful.",
    ]:
        add_line(frame, text, size=11, bullet="•", indent=0.12, spacing=0.95, after=2,
                 align=PP_ALIGN.JUSTIFY)

    _draw_architecture(slide)


def build_feasibility_slide(slide) -> None:
    box = find(slide, "TextBox 8")
    set_autofit_off(box)
    place(box, 0.36, 1.24, 12.58, 4.70)
    clear_text(box)
    frame = box.text_frame

    sections = [
        (
            "Analysis of the feasibility of the idea",
            [
                "Every component is proven technology in isolation; the contribution is the concept-graph spine that connects them. Nothing here needs a research breakthrough.",
                "The prototype already grades a full cohort end to end on commodity hardware — no GPU, no paid API on the critical path, no network access inside the sandbox.",
                "The highest-cost input is roughly two hours of instructor time to author the concept graph once per course. Every insight in the platform traces back to it, and it is reused every semester.",
                "Queue-backed autoscaling workers absorb the spike at a deadline, which is where naive designs fall over.",
            ],
        ),
        (
            "Potential challenges and risks",
            [
                "Executing untrusted code written by capable people, at scale, on a deadline.",
                "An LLM hallucinating an expected output and silently penalising a whole cohort.",
                "Cold start: no labelled data in semester zero.",
                "False plagiarism accusations — the highest-consequence error the system can make.",
                "Faculty adoption: a tool that costs more time than it saves is abandoned regardless of model quality.",
                "Demographic bias in the early-warning model.",
            ],
        ),
        (
            "Strategies for overcoming these challenges",
            [
                "Defence in depth: thirteen isolation layers, one-shot instances, workers holding nothing worth stealing, and the oracle kept outside the guest.",
                "Every generated test executes against the instructor's reference solution before admission; more than ~20% failing halts authoring and flags the brief as ambiguous.",
                "Unsupervised models (clustering, knowledge tracing) ship on day one; the LLM is a teacher that drafts and bootstraps, not a grader. Per-submission LLM cost trends to zero.",
                "Base-code and common-idiom exclusion plus cohort-relative outlier detection; similarity is surfaced as ranked evidence with aligned regions, and faculty decide.",
                "A ten-minute authoring budget and gradebook writeback so faculty never maintain two gradebooks; every override is captured as training data.",
                "A demographic bias audit that blocks deployment when flag rates differ by more than 5% across any protected group, re-run every semester.",
            ],
        ),
    ]
    _render_sections(frame, sections, heading_size=14.5, body_size=11)

    metric_strip(slide, [
        ("< 90 s", "to grade a 24-student, four-lab cohort end to end on a laptop"),
        ("0", "GPUs, paid API calls, or network access needed at grading time"),
        ("~2 h", "of instructor time to author the concept graph, once per course"),
        ("> 20 %", "reference-test failure halts authoring and flags an ambiguous brief"),
    ])


def build_impact_slide(slide) -> None:
    box = find(slide, "TextBox 8")
    set_autofit_off(box)
    place(box, 0.36, 1.24, 12.58, 4.70)
    clear_text(box)
    frame = box.text_frame

    sections = [
        (
            "Potential impact on the target audience",
            [
                "Students — feedback that is evidence rather than a number: \"Empty-input handling 0/8. Test 11 crashed with IndexError at solution.py:14. No length guard found on the input path.\" A ranked next action with the reason, and one-click appeal on any rubric item.",
                "Faculty — a course-health landing view that answers what to re-teach, ranked by how many later concepts depend on it, plus alerts for rubric items that are measuring nothing.",
                "Administrators — live CO–PO attainment with per-student traceability down to individual submissions, and an at-risk view that routes to advising rather than to sanction.",
            ],
        ),
        (
            "Benefits of the solution (social, economic, environmental, etc.)",
            [
                "Social — grading a student can interrogate, contest, and understand. A missing delimiter stops being worth a hundred percent of the marks. Early warning is a support route by design, and its bias audit is mandatory.",
                "Academic — a misconception library and a mastery model that persist across semesters, so what a cohort got wrong in 2026 informs how the topic is taught in 2027.",
                "Economic — faculty minutes per assignment fall semester over semester as auto-release coverage rises; NBA/NAAC attainment reporting collapses from weeks of manual work into a live view.",
                "Operational — deterministic-first design means most of the pipeline is a parser and a test runner, so it runs on ordinary institutional hardware with no GPU and no per-submission API spend.",
                "Equity — mastery is inferred from evidence the student generated, disclosed to them, and auditable; protected attributes are used only to audit the model, never as features.",
            ],
        ),
    ]
    _render_sections(frame, sections, heading_size=14.5, body_size=12)

    # Published targets, not claims: the platform reports itself against these
    # every semester, including when it misses them.
    metric_strip(slide, [
        ("70 %", "auto-release coverage target by semester 2, at under 3% override rate"),
        ("< 3 min", "p95 feedback latency target - feedback loses pedagogical value fast"),
        ("3 weeks", "early-warning lead time before failure, routed to support"),
        ("< 5 %", "maximum early-warning flag-rate gap across any protected group"),
    ])


def build_references_slide(slide) -> None:
    box = find(slide, "TextBox 8")
    set_autofit_off(box)
    place(box, 0.36, 1.24, 12.58, 4.70)
    clear_text(box)
    frame = box.text_frame

    add_line(
        frame, "Details / Links of the reference and research work",
        size=15, bold=True, colour=HEADING_BLUE, underline=True,
        bullet="v", bullet_font="Wingdings", after=7,
    )

    groups = [
        ("Integrity and code similarity", [
            "Schleimer, Wilkerson & Aiken — \"Winnowing: Local Algorithms for Document Fingerprinting\", ACM SIGMOD 2003 (the MOSS algorithm).",
            "Ragkhitwetsagul, Krinke & Clark — \"A comparison of code similarity analysers\", Empirical Software Engineering, 2018.",
        ]),
        ("Knowledge tracing and learning analytics", [
            "Corbett & Anderson — \"Knowledge Tracing: Modeling the Acquisition of Procedural Knowledge\", User Modeling and User-Adapted Interaction, 1994.",
            "Pardos & Heffernan — \"Modeling Individualization in a Bayesian Networks Implementation of Knowledge Tracing\", UMAP 2010.",
            "Ebel & Frisbie — Essentials of Educational Measurement (item difficulty, point-biserial discrimination).",
        ]),
        ("Program analysis and partial credit", [
            "Pawlik & Augsten — \"Tree edit distance: Robust and memory-efficient\" (APTED), Information Systems, 2016.",
            "Gulwani, Radiček & Zuleger — \"Automated Clustering and Program Repair for Introductory Programming Assignments\", PLDI 2018.",
            "tree-sitter — incremental, error-tolerant parsing library, https://tree-sitter.github.io",
        ]),
        ("Clustering, NLI, and sandboxing", [
            "Campello, Moulavi & Sander — \"Density-Based Clustering Based on Hierarchical Density Estimates\" (HDBSCAN), PAKDD 2013.",
            "He, Gao & Chen — \"DeBERTaV3: Improving DeBERTa using ELECTRA-Style Pre-Training\", ICLR 2023.",
            "Agache et al. — \"Firecracker: Lightweight Virtualization for Serverless Applications\", USENIX NSDI 2020.",
        ]),
        ("Standards and institutional context", [
            "1EdTech (IMS Global) — Learning Tools Interoperability 1.3 and Assignment and Grade Services, https://www.imsglobal.org/spec/lti/v1p3",
            "National Board of Accreditation, India — CO–PO attainment guidelines for UG engineering programmes.",
            "Digital Personal Data Protection Act, 2023 (India) — treatment of inferred attributes as personal data.",
        ]),
    ]
    for heading, entries in groups:
        add_line(frame, heading, size=12, bold=True, colour=ACCENT_BLUE,
                 bullet="•", indent=0.0, spacing=0.95, before=3, after=1)
        for entry in entries:
            add_line(frame, entry, size=10, bullet="–", indent=0.30, spacing=0.95, after=1)
    drop_empty_paragraphs(frame)

    footer = slide.shapes.add_textbox(Inches(0.34), Inches(6.08), Inches(12.66), Inches(0.62))
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(0xEC, 0xF2, 0xF9)
    footer.line.color.rgb = LAYER_FILL_ALT
    footer.line.width = Pt(0.75)
    set_autofit_off(footer)
    clear_text(footer)
    footer.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_line(
        footer.text_frame,
        "Working prototype and full architecture specification: github.com/Varshith-06/evalpro",
        size=12, bold=True, colour=LAYER_FILL, bullet=None, align=PP_ALIGN.CENTER, first=True, spacing=0.95,
    )
    add_line(
        footer.text_frame,
        "The demo course builds itself on first run - every score, mastery estimate and attainment figure "
        "in it is produced by the real cascade, not seeded.",
        size=9, colour=BODY_BLACK, bullet=None, align=PP_ALIGN.CENTER, spacing=0.92,
    )


def metric_strip(slide, items: list[tuple[str, str]], top: float = 6.02, height: float = 0.76) -> None:
    """A row of measured figures along the foot of a slide.

    Judges read numbers before prose, and the numbers here are measured on the
    working prototype rather than projected, so they are given the space to be
    read on their own.
    """
    left = Inches(0.34)
    width = Inches(2.98)
    gap = Inches(0.22)
    for index, (figure, caption) in enumerate(items):
        x = Emu(int(left) + index * (int(width) + int(gap)))
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(top), width, Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xEC, 0xF2, 0xF9)
        box.line.color.rgb = LAYER_FILL_ALT
        box.line.width = Pt(0.75)
        box.shadow.inherit = False
        set_autofit_off(box)
        clear_text(box)
        frame = box.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_line(frame, figure, size=15, bold=True, colour=LAYER_FILL, bullet=None,
                 align=PP_ALIGN.CENTER, first=True, spacing=0.9)
        add_line(frame, caption, size=8, colour=BODY_BLACK, bullet=None,
                 align=PP_ALIGN.CENTER, spacing=0.88)


def _render_sections(frame, sections, heading_size: float, body_size: float) -> None:
    """Pointer heading (preserved verbatim from the template) then the content."""
    for heading, bullets in sections:
        add_line(
            frame, heading, size=heading_size, bold=True, colour=HEADING_BLUE,
            underline=True, bullet="v", bullet_font="Wingdings",
            before=7, after=3, align=PP_ALIGN.LEFT,
        )
        for text in bullets:
            add_line(
                frame, text, size=body_size, bullet="•", indent=0.12,
                spacing=0.95, after=2, align=PP_ALIGN.JUSTIFY,
            )


# ==========================================================================
# Architecture diagram, drawn as native shapes
# ==========================================================================
LAYERS = [
    ("L0  INSTITUTIONAL CONTEXT", "LTI 1.3 / SIS sync\nroster · course outcomes\nconcept graph"),
    ("L1  SENSING", "authoring → cascade\nsandbox → gate\nper-rubric-item evidence"),
    ("L2  ACCUMULATION", "knowledge tracing\nitem analysis\nmisconception clusters"),
    ("L3  ACTION", "remediation · re-teach\nearly warning\nCO–PO attainment"),
]


def _draw_architecture(slide) -> None:
    top = Inches(1.22)
    height = Inches(1.50)
    left = Inches(0.36)
    width = Inches(2.86)
    gap = Inches(0.32)

    for index, (title, body) in enumerate(LAYERS):
        x = Emu(int(left) + index * (int(width) + int(gap)))
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = LAYER_FILL if index % 2 == 0 else LAYER_FILL_ALT
        box.line.color.rgb = WHITE
        box.line.width = Pt(1)
        box.shadow.inherit = False

        frame = box.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_autofit_off(box)
        clear_text(box)
        add_line(frame, title, size=10.5, bold=True, colour=WHITE,
                 bullet=None, align=PP_ALIGN.CENTER, first=True, spacing=0.95, after=3)
        add_line(frame, body, size=8.5, colour=WHITE, bullet=None,
                 align=PP_ALIGN.CENTER, spacing=0.9)

        if index < len(LAYERS) - 1:
            arrow_left = Emu(int(x) + int(width))
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, arrow_left, Emu(int(top) + int(height) // 2),
                Emu(int(arrow_left) + int(gap)), Emu(int(top) + int(height) // 2),
            )
            connector.line.color.rgb = HEADING_BLUE
            connector.line.width = Pt(2.25)
            line = connector.line._get_or_add_ln()
            tail = line.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
            line.append(tail)

    caption = slide.shapes.add_textbox(left, Emu(int(top) + int(height) + int(Inches(0.06))),
                                       Inches(12.55), Inches(0.44))
    set_autofit_off(caption)
    clear_text(caption)
    add_line(
        caption.text_frame,
        "The single field that connects L1 to L2 is RubricItem.concept_ids — without it you have twelve "
        "disconnected gradebooks; with it, every submission is evidence about a named competency.",
        size=9, bold=True, colour=HEADING_BLUE, bullet=None, first=True,
        align=PP_ALIGN.CENTER, spacing=0.95,
    )


# ==========================================================================
# Entry point
# ==========================================================================
def main() -> int:
    if not TEMPLATE.exists():
        print(f"template not found: {TEMPLATE}", file=sys.stderr)
        return 1
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(TEMPLATE, OUTPUT)

    prs = Presentation(str(OUTPUT))
    slides = list(prs.slides)

    build_title_slide(slides[0])
    build_solution_slide(slides[1])
    build_technical_slide(slides[2])
    build_feasibility_slide(slides[3])
    build_impact_slide(slides[4])
    build_references_slide(slides[5])

    # The template's own instructions slide says the deck must be six slides
    # including the title page, and that this slide may be deleted before upload.
    delete_slide(prs, 6)

    _set_team_name(prs)
    prs.save(str(OUTPUT))

    print(f"wrote {OUTPUT.relative_to(REPO)} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
    return 0


def _set_team_name(prs) -> None:
    """The oval in the top-left of every content slide reads 'Your Team Name'."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip() == "Your Team Name":
                        run.text = TEAM["team_name"] if not TEAM["team_name"].startswith("[") else "Your Team Name"


if __name__ == "__main__":
    raise SystemExit(main())
